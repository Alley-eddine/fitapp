"""Generate the Bloc 1 soutenance deck (FitCoach AI) as a styled .pptx."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- palette (FitCoach AI : dark slate + cyan) ----
BG     = RGBColor(0x0F, 0x17, 0x2A)
PANEL  = RGBColor(0x1E, 0x29, 0x3B)
PANEL2 = RGBColor(0x27, 0x34, 0x49)
CYAN   = RGBColor(0x22, 0xD3, 0xEE)
CYAN_D = RGBColor(0x0E, 0x74, 0x90)
WHITE  = RGBColor(0xF8, 0xFA, 0xFC)
MUTED  = RGBColor(0x94, 0xA3, 0xB8)
GREEN  = RGBColor(0x34, 0xD3, 0x99)
RED    = RGBColor(0xF8, 0x71, 0x71)
AMBER  = RGBColor(0xFB, 0xBF, 0x24)
DARK   = RGBColor(0x0B, 0x12, 0x20)

HF = "Trebuchet MS"   # headers
BF = "Calibri"        # body

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height


def _fill(shape, color):
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()


def rect(slide, x, y, w, h, color, rounded=False):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        x, y, w, h)
    _fill(shp, color)
    return shp


def text(slide, x, y, w, h, runs, size=16, color=WHITE, bold=False, font=BF,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing=1.0):
    tb = slide.shapes.add_textbox(x, y, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    items = runs if isinstance(runs, list) else [(runs, {})]
    first = True
    for content, opts in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = opts.get("align", align)
        p.line_spacing = opts.get("line_spacing", line_spacing)
        p.space_after = Pt(opts.get("space_after", 4))
        p.space_before = Pt(opts.get("space_before", 0))
        sub = opts.get("parts", [(content, {})])
        for t2, o2 in sub:
            r = p.add_run(); r.text = t2
            r.font.size = Pt(o2.get("size", size))
            r.font.bold = o2.get("bold", bold)
            r.font.name = o2.get("font", font)
            r.font.color.rgb = o2.get("color", color)
    return tb


def base(eyebrow, title, title_size=33):
    s = prs.slides.add_slide(BLANK)
    rect(s, 0, 0, SW, SH, BG)
    rect(s, Inches(0.7), Inches(0.62), Inches(0.16), Inches(0.16), CYAN)  # motif
    text(s, Inches(0.95), Inches(0.55), Inches(11.6), Inches(0.3),
         eyebrow.upper(), size=12.5, color=CYAN, bold=True, font=HF)
    text(s, Inches(0.7), Inches(0.92), Inches(11.9), Inches(0.9),
         title, size=title_size, color=WHITE, bold=True, font=HF, line_spacing=1.0)
    return s


def bullets(slide, x, y, w, h, items, size=15, gap=7, color=WHITE):
    runs = []
    for it in items:
        if isinstance(it, tuple):
            label, val = it
            runs.append((None, {"space_after": gap, "parts": [
                ("▸  ", {"color": CYAN, "bold": True, "size": size}),
                (label, {"color": color, "bold": True, "size": size}),
                (val, {"color": MUTED, "size": size})]}))
        else:
            runs.append((None, {"space_after": gap, "parts": [
                ("▸  ", {"color": CYAN, "bold": True, "size": size}),
                (it, {"color": color, "size": size})]}))
    text(slide, x, y, w, h, runs)


def table(slide, x, y, w, headers, rows, col_w=None, fs=11, header_fs=11, row_h=0.3):
    nrows, ncols = len(rows) + 1, len(headers)
    gtbl = slide.shapes.add_table(nrows, ncols, x, y, w, Inches(row_h * nrows)).table
    gtbl.first_row = False; gtbl.horz_banding = False
    # remove default style banding by setting our own fills
    if col_w:
        total = sum(col_w)
        for i, cw in enumerate(col_w):
            gtbl.columns[i].width = int(w * cw / total)
    for j, htxt in enumerate(headers):
        c = gtbl.cell(0, j)
        c.fill.solid(); c.fill.fore_color.rgb = CYAN_D
        c.margin_left = Pt(6); c.margin_right = Pt(6); c.margin_top = Pt(3); c.margin_bottom = Pt(3)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = c.text_frame.paragraphs[0]; r = p.add_run(); r.text = htxt
        r.font.size = Pt(header_fs); r.font.bold = True; r.font.name = HF; r.font.color.rgb = WHITE
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            c = gtbl.cell(i + 1, j)
            c.fill.solid(); c.fill.fore_color.rgb = PANEL if i % 2 == 0 else PANEL2
            c.margin_left = Pt(6); c.margin_right = Pt(6); c.margin_top = Pt(2); c.margin_bottom = Pt(2)
            c.vertical_anchor = MSO_ANCHOR.MIDDLE
            p = c.text_frame.paragraphs[0]
            # allow inline color via tuple (text, color)
            if isinstance(val, tuple):
                txt, col = val; bold = True
            else:
                txt, col, bold = val, WHITE, False
            r = p.add_run(); r.text = txt
            r.font.size = Pt(fs); r.font.name = BF; r.font.color.rgb = col; r.font.bold = bold
    return gtbl


def card(slide, x, y, w, h, title, items, accent=CYAN):
    c = rect(slide, x, y, w, h, PANEL, rounded=True)
    text(slide, x + Inches(0.18), y + Inches(0.12), w - Inches(0.36), Inches(0.4),
         title, size=14.5, color=accent, bold=True, font=HF)
    runs = [(None, {"space_after": 4, "parts": [("• ", {"color": accent, "size": 12}),
            (it, {"color": WHITE, "size": 12})]}) for it in items]
    text(slide, x + Inches(0.18), y + Inches(0.55), w - Inches(0.36), h - Inches(0.7), runs)


# ============================ SLIDE 1 — TITLE ============================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, BG)
rect(s, 0, 0, Inches(0.35), SH, CYAN)  # left bar motif
text(s, Inches(1.0), Inches(2.0), Inches(11.5), Inches(0.4),
     "RNCP 39583 — EXPERT EN DÉVELOPPEMENT LOGICIEL", size=14, color=CYAN, bold=True, font=HF)
text(s, Inches(0.95), Inches(2.5), Inches(11.6), Inches(1.4),
     [(None, {"parts": [("FitCoach ", {"color": WHITE, "bold": True, "size": 60, "font": HF}),
                         ("AI", {"color": CYAN, "bold": True, "size": 60, "font": HF})]})])
text(s, Inches(1.0), Inches(3.7), Inches(11.0), Inches(0.9),
     "Cadrer un projet de développement d'applications logicielles — Bloc 1",
     size=22, color=WHITE, bold=True, font=HF)
text(s, Inches(1.0), Inches(4.5), Inches(11.0), Inches(0.8),
     "Plateforme de coaching sportif : transmission coach ↔ élève, tracée et sans saisie.",
     size=16, color=MUTED)
text(s, Inches(1.0), Inches(6.4), Inches(11.0), Inches(0.5),
     [(None, {"parts": [("Commanditaire : ", {"color": MUTED, "size": 14}),
                         ("Markus, coach sportif privé", {"color": WHITE, "bold": True, "size": 14}),
                         ("    ·    Candidat : Alley Eddine", {"color": MUTED, "size": 14})]})])

# ============================ SLIDE 2 — SOMMAIRE ============================
s = base("Plan de la présentation", "Sommaire")
left = ["Cartographie des parties prenantes — C1.1.1",
        "Analyse de la demande — C1.1.2",
        "SWOT & concurrence — C1.2.1",
        "Faisabilité & audit de l'existant — C1.2.2",
        "Cartographie des risques — C1.2.3",
        "Veille technique & réglementaire — C1.3.1"]
right = ["Architecture technique (comparatif) — C1.3.2",
         "Arbitrages techniques — C1.3.2",
         "Charge de travail — C1.4.1",
         "Coût & budget — C1.4.2",
         "Modélisation de l'architecture — C1.5",
         "Décisions & axes de solutions — C1.6"]
for i, (col, x) in enumerate([(left, Inches(0.9)), (right, Inches(7.0))]):
    runs = []
    for k, it in enumerate(col):
        n = i * 6 + k + 1
        runs.append((None, {"space_after": 12, "parts": [
            (f"{n:02d}  ", {"color": CYAN, "bold": True, "size": 17, "font": HF}),
            (it, {"color": WHITE, "size": 15})]}))
    text(s, x, Inches(2.0), Inches(5.6), Inches(4.8), runs)

# ============================ SLIDE 3 — PARTIES PRENANTES ============================
s = base("C1.1.1 — Cartographie des parties prenantes", "Les acteurs du projet")
table(s, Inches(0.7), Inches(1.9), Inches(11.9),
      ["Acteur", "Rôle", "Implication", "Canal"],
      [[("Markus — commanditaire", CYAN), "Décideur · expert métier · futur client", "Élevée", "Entretiens"],
       ["Développeur (moi)", "Conception & développement complet", "Permanente", "—"],
       ["Coach (rôle app)", "Crée programmes/nutrition, suit ses élèves", "Quotidienne", "Desktop"],
       ["Élève (rôle app)", "Exécute séances pré-réglées, renseigne son suivi", "Quotidienne", "PWA mobile"],
       ["Utilisateur lambda", "Self-tracking sans coach", "Variable", "PWA"],
       ["DPO / conformité", "RGPD — données de santé (poids, mensurations)", "Ponctuelle", "Audits"],
       ["Hébergeur", "Disponibilité de la plateforme", "Continue", "Infra"]],
      col_w=[3, 5.2, 2, 2], fs=11.5, row_h=0.62)
text(s, Inches(0.7), Inches(6.85), Inches(11.9), Inches(0.4),
     [(None, {"parts": [("Supports : ", {"color": CYAN, "bold": True, "size": 12}),
        ("Coach → desktop (vue d'ensemble) · Élève → mobile (séance du jour) · Lambda → PWA standard.",
         {"color": MUTED, "size": 12})]})])

# ============================ SLIDE 4 — DEMANDE ============================
s = base("C1.1.2 — Analyse de la demande", "Du suivi WhatsApp à une plateforme tracée")
text(s, Inches(0.7), Inches(1.85), Inches(11.9), Inches(0.6),
     [(None, {"parts": [("Existant : ", {"color": CYAN, "bold": True, "size": 14}),
        ("Markus me coache depuis ~1,5 mois — programmes & repas envoyés par email (PDF), suivi sur WhatsApp.",
         {"color": WHITE, "size": 14})]})])
card(s, Inches(0.7), Inches(2.55), Inches(5.8), Inches(2.5), "Besoins identifiés",
     ["Centraliser la transmission (fini les PDF par mail)",
      "Structurer et tracer le suivi (WhatsApp = pas d'historique)",
      "Élève : ne rien régler — séance & repas déjà prêts",
      "Remontée automatique des données au coach"])
card(s, Inches(6.8), Inches(2.55), Inches(5.8), Inches(2.5), "Pistes de solution", [
      "Programmes versionnés par phase, assignés à l'élève",
      "Séances pré-réglées · repas imposés + variété par IA",
      "Notifications push contextuelles « voix du coach »",
      "Suivi poids / mensurations en graphiques"], accent=GREEN)
text(s, Inches(0.7), Inches(5.35), Inches(11.9), Inches(1.3),
     [(None, {"parts": [("Problématique : ", {"color": CYAN, "bold": True, "size": 15}),
        ("comment centraliser et tracer la transmission coach ↔ élève, en supprimant l'éparpillement "
         "email / WhatsApp, pour un suivi fiable et un minimum de saisie ?",
         {"color": WHITE, "size": 15})]})], )

# ============================ SLIDE 5 — SWOT ============================
s = base("C1.2.1 — Analyse SWOT", "Forces, faiblesses, opportunités, menaces")
gx, gy, cw, ch = Inches(0.7), Inches(1.9), Inches(5.85), Inches(2.3)
card(s, gx, gy, cw, ch, "Forces", [
    "Besoin réel validé par mon propre coach",
    "Full-stack + DevOps (microservices, observabilité)",
    "Différenciation : nutrition cadrée + IA, zéro-config"], accent=GREEN)
card(s, Inches(6.75), gy, cw, ch, "Faiblesses", [
    "Ressource unique (dev solo — bus factor)",
    "Périmètre large (entraînement + nutrition + paiement)",
    "Adoption élèves à prouver"], accent=AMBER)
card(s, gx, Inches(4.35), cw, ch, "Opportunités", [
    "Coachs indépendants bricolent (email/WhatsApp/PDF)",
    "Modèle B2B2C : abonnement coach (récurrent)",
    "Extension multi-coachs / communauté"], accent=CYAN)
card(s, Inches(6.75), Inches(4.35), cw, ch, "Menaces", [
    "Données de santé sensibles → RGPD / cyber",
    "Concurrents établis (Trainerize, TrueCoach)",
    "Délais (solo + certification)"], accent=RED)
text(s, Inches(0.7), Inches(6.8), Inches(11.9), Inches(0.4),
     [(None, {"parts": [("Impact environnemental : ", {"color": CYAN, "bold": True, "size": 12}),
        ("conteneurs scalés à la demande, hébergeur éco-responsable visé.", {"color": MUTED, "size": 12})]})])

# ============================ SLIDE 6 — CONCURRENCE ============================
s = base("C1.2.1 — Concurrence & positionnement", "Un marché servi… mais pas pour ce coach")
card(s, Inches(0.7), Inches(1.9), Inches(5.85), Inches(3.4), "La concurrence", [
    "Trainerize / TrueCoach : coach-élève complet, mais US, cher, nutrition générique",
    "MyFitnessPal / Freeletics : grand public, sans relation coach",
    "La réalité de Markus : email + WhatsApp + PDF → zéro structure, zéro traçabilité"], accent=MUTED)
card(s, Inches(6.75), Inches(1.9), Inches(5.85), Inches(3.4), "Mes différenciateurs", [
    "Séance zéro-config : tout est déjà réglé",
    "Nutrition imposée par le coach + variété par IA",
    "Remontée automatique des données au coach",
    "Notifications « voix du coach » contextuelles",
    "Phases qui matérialisent le suivi dans le temps"], accent=CYAN)
text(s, Inches(0.7), Inches(5.6), Inches(11.9), Inches(1.0),
     [(None, {"parts": [("Positionnement : ", {"color": CYAN, "bold": True, "size": 15}),
        ("non pas un tracker fitness de plus, mais l'outil de transmission et de suivi du coach "
         "indépendant — pensé France et zéro friction.", {"color": WHITE, "size": 15})]})])

# ============================ SLIDE 7 — FAISABILITÉ ============================
s = base("C1.2.2 — Faisabilité & audit de l'existant", "Verdict : faisable")
card(s, Inches(0.7), Inches(1.9), Inches(5.85), Inches(3.3), "Audit de l'existant", [
    "Email + PDF : info dispersée, pas de remontée",
    "WhatsApp : suivi non structuré, pas de preuve d'adhérence",
    "Outils marché : self-service sans coach, ou US lourds/chers",
    "→ Aucun ne couvre transmission tracée + zéro-config + nutrition IA"], accent=AMBER)
card(s, Inches(6.75), Inches(1.9), Inches(5.85), Inches(3.3), "Stack cible", [
    "Front : Next.js 15 en PWA (installable, web-push)",
    "Back : 5 microservices Fastify (TypeScript)",
    "Data : PostgreSQL · Auth : JWT",
    "Stripe · Groq · Resend / Twilio / web-push",
    "Docker · Prometheus + Grafana"], accent=CYAN)
text(s, Inches(0.7), Inches(5.5), Inches(11.9), Inches(1.0),
     [(None, {"parts": [("✅ Verdict : ", {"color": GREEN, "bold": True, "size": 15}),
        ("faisable — stack maîtrisée, MVP cadré, socle déjà amorcé. Contraintes : solo, données de "
         "santé (RGPD), multi-rôles, fluidité.", {"color": WHITE, "size": 15})]})])

# ============================ SLIDE 8 — RISQUES ============================
s = base("C1.2.3 — Cartographie des risques", "Criticité = probabilité × impact")
table(s, Inches(0.7), Inches(1.9), Inches(11.9),
      ["Risque", "Crit.", "Plan d'action"],
      [["Fuite de données de santé", ("Critique", RED), "JWT, clé interne, Zod, TLS, minimisation"],
       ["Indisponibilité d'un service", ("Critique", RED), "Découplage, health checks, monitoring + alertes"],
       ["Retard de livraison (solo)", ("Majeur", AMBER), "MVP priorisé, Kanban, jalons RNCP"],
       ["Perte de données", ("Majeur", AMBER), "Backups PostgreSQL, migrations idempotentes"],
       ["Non-conformité RGPD (santé)", ("Majeur", AMBER), "RGPD by design, minimisation, durée limitée"],
       ["Adoption faible des élèves", ("Modéré", GREEN), "UX zéro-config, notifications « voix du coach »"]],
      col_w=[3.4, 1.4, 6.2], fs=11.5, row_h=0.6)
text(s, Inches(0.7), Inches(6.75), Inches(11.9), Inches(0.5),
     [(None, {"parts": [("Indicateurs de contrôle : ", {"color": CYAN, "bold": True, "size": 12}),
        ("uptime par service · taux d'erreur · couverture de tests · latence — via Prometheus / Grafana.",
         {"color": MUTED, "size": 12})]})])

# ============================ SLIDE 9 — VEILLE ============================
s = base("C1.3.1 — Veille technique & réglementaire", "Anticiper failles & conformité")
card(s, Inches(0.7), Inches(1.9), Inches(5.85), Inches(3.0), "Veille technique", [
    "Next.js / Fastify : releases & changelogs",
    "Node / npm : security advisories · Dependabot",
    "OWASP : Top 10, bonnes pratiques",
    "Stripe / Groq : changelogs API"], accent=CYAN)
card(s, Inches(6.75), Inches(1.9), Inches(5.85), Inches(3.0), "Veille réglementaire", [
    "CNIL : RGPD, données de santé (catégorie particulière)",
    "ANSSI : sécurité des systèmes"], accent=GREEN)
text(s, Inches(0.7), Inches(5.2), Inches(11.9), Inches(1.0),
     [(None, {"parts": [("Outils : ", {"color": CYAN, "bold": True, "size": 14}),
        ("Feedly · GitHub Watch · Dependabot · Google Alerts.   ", {"color": WHITE, "size": 14}),
        ("Fréquences : ", {"color": CYAN, "bold": True, "size": 14}),
        ("deps → hebdo · CNIL/ANSSI → mensuel.", {"color": WHITE, "size": 14})]})])

# ============================ SLIDE 10 — ARCHI (ÉLIMINATOIRE) ============================
s = base("C1.3.2 — Étude comparative (compétence éliminatoire)", "Choix : microservices")
table(s, Inches(0.7), Inches(1.85), Inches(11.9),
      ["Critère", "BaaS (Supabase)", "Monolithe", "Microservices ✓"],
      [["Maîtrise & sécurité", "Moyen (boîte noire)", "Élevé", ("Élevé (isolation)", CYAN)],
       ["Séparation des responsabilités", "Faible", "Moyen", ("Excellente", CYAN)],
       ["Scalabilité indépendante", "Limitée", "Faible", ("Excellente (IA)", CYAN)],
       ["Résilience (panne isolée)", "Moyenne", "Faible", ("Bonne", CYAN)],
       ["Démonstration niveau Expert", "Faible", "Moyen", ("Forte", CYAN)],
       ["Complexité opérationnelle", "Faible", "Faible", ("Élevée → observabilité", AMBER)]],
      col_w=[3.2, 2.9, 2, 3.0], fs=11, row_h=0.52)
text(s, Inches(0.7), Inches(5.95), Inches(11.9), Inches(1.3),
     [(None, {"parts": [("Retenu : ", {"color": CYAN, "bold": True, "size": 13}),
        ("auth / paiement / IA / notifications ont des profils de sécurité et de scaling distincts → "
         "déploiement & scaling indépendants, isolation des pannes, maîtrise du distribué (niveau Expert). "
         "Couplage maîtrisé (JWT + clé interne), complexité compensée par l'observabilité.",
         {"color": WHITE, "size": 13})]})])

# ============================ SLIDE 11 — ARBITRAGES ============================
s = base("C1.3.2 — Arbitrages techniques", "Retenu vs écarté")
table(s, Inches(0.7), Inches(1.9), Inches(11.9),
      ["Option", "Décision", "Justification"],
      [["Microservices Fastify", ("✅ RETENU", GREEN), "Domaines hétérogènes, scaling/sécurité indépendants"],
       ["Next.js 15 en PWA", ("✅ RETENU", GREEN), "1 codebase installable (coach desktop + élève mobile), web-push"],
       ["PostgreSQL + SQL direct", ("✅ RETENU", GREEN), "Maîtrise des requêtes, perf, démontre la compétence BDD"],
       ["JWT + clé interne", ("✅ RETENU", GREEN), "Auth maîtrisée + isolation des appels internes"],
       ["Zod partagé front/back", ("✅ RETENU", GREEN), "Une seule source de validation"],
       ["BaaS (Supabase)", ("❌ ÉCARTÉ", RED), "Sécurité en boîte noire, faible démonstration"],
       ["App native (2 apps)", ("❌ ÉCARTÉ", RED), "PWA couvre offline/push sans double maintenance"]],
      col_w=[3.0, 2.0, 6.6], fs=11, row_h=0.55)

# ============================ SLIDE 12 — CHARGE ============================
s = base("C1.4.1 — Charge de travail", "Estimation en jours-homme")
table(s, Inches(0.7), Inches(1.9), Inches(8.4),
      ["Phase", "j/h"],
      [["Cadrage & architecture", "5"], ["Auth & rôles (coach/élève/lambda)", "8"],
       ["Service API métier", "12"], ["Service IA (nutrition cadrée)", "6"],
       ["Service Paiement (Stripe)", "6"], ["Service Notifications (mail/SMS/push)", "6"],
       ["Front Next.js PWA (3 rôles)", "14"], ["Tests & CI/CD", "6"],
       ["Monitoring & déploiement", "5"], ["Documentation & recette", "5"]],
      col_w=[7, 1], fs=11.5, row_h=0.42)
rect(s, Inches(9.5), Inches(1.9), Inches(3.1), Inches(1.7), PANEL, rounded=True)
text(s, Inches(9.5), Inches(2.05), Inches(3.1), Inches(0.6), "TOTAL", size=13, color=MUTED, bold=True, align=PP_ALIGN.CENTER, font=HF)
text(s, Inches(9.5), Inches(2.45), Inches(3.1), Inches(0.9), "≈ 73 j/h", size=34, color=CYAN, bold=True, align=PP_ALIGN.CENTER, font=HF)
text(s, Inches(9.5), Inches(3.25), Inches(3.1), Inches(0.4), "≈ 14-15 semaines", size=12, color=MUTED, align=PP_ALIGN.CENTER)
card(s, Inches(9.5), Inches(3.9), Inches(3.1), Inches(2.7), "Hiérarchie", [
    "Principales : auth, programmes, séance, suivi",
    "Secondaires : nutrition IA, paiement, notifs",
    "Complémentaires : graphiques, communauté"], accent=CYAN)

# ============================ SLIDE 13 — COÛT ============================
s = base("C1.4.2 — Coût & budget", "Coûts réels vs valeur marché")
card(s, Inches(0.7), Inches(1.95), Inches(5.85), Inches(3.2), "Coûts directs (solo)", [
    "Hébergement (~5 mois) : ~150 €",
    "Domaine + SSL : ~20 €",
    "APIs (Stripe test, Twilio, Resend, Groq) : ~80 €",
    "Outils & monitoring : ~30 €"], accent=CYAN)
rect(s, Inches(6.75), Inches(1.95), Inches(5.85), Inches(3.2), PANEL, rounded=True)
text(s, Inches(6.95), Inches(2.1), Inches(5.5), Inches(0.4), "VALEUR MARCHÉ (RÉFÉRENCE)", size=13, color=GREEN, bold=True, font=HF)
text(s, Inches(6.95), Inches(2.6), Inches(5.5), Inches(0.5), "Dev senior 500 €/j × 73 j", size=14, color=WHITE)
text(s, Inches(6.95), Inches(3.05), Inches(5.5), Inches(1.0), "≈ 37 000 €", size=40, color=GREEN, bold=True, font=HF)
text(s, Inches(6.95), Inches(4.2), Inches(5.5), Inches(0.7), "+ infra/services annuels ~600 €", size=13, color=MUTED)
# big direct cost callout
rect(s, Inches(0.7), Inches(5.4), Inches(5.85), Inches(1.1), PANEL2, rounded=True)
text(s, Inches(0.9), Inches(5.6), Inches(5.5), Inches(0.7),
     [(None, {"parts": [("Total coûts directs ≈ ", {"color": WHITE, "size": 16}),
                         ("280 €", {"color": CYAN, "bold": True, "size": 22, "font": HF})]})])
text(s, Inches(6.75), Inches(5.5), Inches(5.85), Inches(1.0),
     "L'écart coût réel / valeur illustre le ROI pour un coach qui s'abonne plutôt que de faire développer.",
     size=13, color=MUTED)

# ============================ SLIDE 14 — MODÉLISATION ============================
s = base("C1.5 — Modélisation de l'architecture", "UML composants + ERD")
# front box
def box(x, y, w, h, label, color=PANEL, tcol=WHITE, fs=12):
    rect(s, x, y, w, h, color, rounded=True)
    lines = label.split("\n")
    runs = [(None, {"align": PP_ALIGN.CENTER, "space_after": 0,
                    "parts": [(ln, {"size": fs, "bold": True, "font": HF, "color": tcol})]})
            for ln in lines]
    text(s, x, y, w, h, runs, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
box(Inches(3.9), Inches(1.85), Inches(5.5), Inches(0.8), "Front Next.js 15 — PWA  (coach · élève · lambda)", CYAN_D, WHITE, 13)
svc = ["auth\n3001", "api\n3002", "ai\n3003", "notifications\n3004", "payment\n3005"]
for i, name in enumerate(svc):
    box(Inches(0.7 + i * 2.45), Inches(3.1), Inches(2.2), Inches(0.95), name, PANEL, CYAN, 12)
box(Inches(3.0), Inches(4.5), Inches(3.2), Inches(0.8), "PostgreSQL", PANEL2, WHITE, 13)
box(Inches(7.1), Inches(4.5), Inches(3.2), Inches(0.8), "Prometheus + Grafana", PANEL2, WHITE, 12)
text(s, Inches(0.7), Inches(5.55), Inches(11.9), Inches(1.4),
     [(None, {"space_after": 5, "parts": [("Entités clés : ", {"color": CYAN, "bold": True, "size": 13}),
        ("users (role) · coach_students · profiles · measurements · training_programs (phases) · "
         "program_assignments · workouts/exercises · nutrition_plans · recipes · notification_logs.",
         {"color": WHITE, "size": 13})]}),
      (None, {"parts": [("Formalisme : ", {"color": CYAN, "bold": True, "size": 13}),
        ("UML (composants) + ERD (données).   Sécurisé (JWT + clé interne + Zod), observable, évolutif.",
         {"color": MUTED, "size": 13})]})])

# ============================ SLIDE 15 — DÉCISIONS ============================
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, SW, SH, BG)
rect(s, 0, 0, Inches(0.35), SH, CYAN)
text(s, Inches(1.0), Inches(0.7), Inches(11.5), Inches(0.4),
     "C1.6 — DÉCISIONS & AXES DE SOLUTIONS", size=13, color=CYAN, bold=True, font=HF)
text(s, Inches(1.0), Inches(1.15), Inches(11.5), Inches(0.8), "Synthèse", size=33, color=WHITE, bold=True, font=HF)
bullets(s, Inches(1.0), Inches(2.2), Inches(11.3), Inches(3.2), [
    "Plateforme coach ↔ élève centralisée, multi-rôles, tracée",
    "Séances zéro-config + nutrition imposée par le coach, variée par IA",
    "Suivi automatique (poids, mensurations, adhérence) en graphiques",
    "Notifications push contextuelles « voix du coach »",
    "Microservices + Next.js PWA, sécurité & RGPD by design, observabilité",
    "Modèle B2B2C : abonnement coach + offre self-service (lambda)"], size=16, gap=10)
rect(s, Inches(1.0), Inches(5.55), Inches(11.3), Inches(1.2), PANEL, rounded=True)
text(s, Inches(1.25), Inches(5.75), Inches(10.8), Inches(0.8),
     [(None, {"parts": [("Pour Markus : ", {"color": CYAN, "bold": True, "size": 15}),
        ("« Tes élèves ouvrent l'app, tout est déjà réglé, et toi tu vois leur progression en temps réel "
         "— sans une minute de saisie en plus. »", {"color": WHITE, "size": 15})]})],
     anchor=MSO_ANCHOR.MIDDLE)

out = "docs/rncp/bloc1/Bloc1_FitCoachAI.pptx"
prs.save(out)
print("OK ->", out, "|", len(prs.slides.__iter__.__self__._sldIdLst), "slides")
